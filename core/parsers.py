"""parsers.py — 多格式文档解析，返回文本 chunk 列表

5 个解析轨道（对应需求文档 2.2 节）:
  Track 1 structured_doc : .dita / .ditamap / .xml / .md / .txt  → 章节级切块
  Track 2 rich_text       : .pdf / .docx / .pptx / .ppt           → 段落/页面块
  Track 3 table           : .xlsx / .xls / .csv                   → 行列转文本段
  Track 4 log             : .log / .json                          → 字段抽取/事件聚合
  Track 5 restricted      : 邮件/聊天/草稿                        → 不直接入库，需人工审核
"""
from __future__ import annotations

import json
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import List

from .config import (
    PARSE_TRACK_LOG,
    PARSE_TRACK_RICH_TEXT,
    PARSE_TRACK_STRUCTURED_DOC,
    PARSE_TRACK_TABLE,
)

# ── 各轨道支持的文件后缀 ──────────────────────────────────────────────────────

SUPPORTED_SUFFIXES = {
    # Track 1
    ".md", ".txt", ".dita", ".ditamap", ".xml",
    # Track 2
    ".pdf", ".ppt", ".pptx", ".docx",
    # Track 3
    ".xlsx", ".xls", ".csv",
    # Track 4
    ".log", ".json",
}

# ── 后缀 → 解析轨道映射 ───────────────────────────────────────────────────────

_TRACK_MAP: dict[str, str] = {
    ".md":      PARSE_TRACK_STRUCTURED_DOC,
    ".txt":     PARSE_TRACK_STRUCTURED_DOC,
    ".dita":    PARSE_TRACK_STRUCTURED_DOC,
    ".ditamap": PARSE_TRACK_STRUCTURED_DOC,
    ".xml":     PARSE_TRACK_STRUCTURED_DOC,
    ".pdf":     PARSE_TRACK_RICH_TEXT,
    ".docx":    PARSE_TRACK_RICH_TEXT,
    ".pptx":    PARSE_TRACK_RICH_TEXT,
    ".ppt":     PARSE_TRACK_RICH_TEXT,
    ".xlsx":    PARSE_TRACK_TABLE,
    ".xls":     PARSE_TRACK_TABLE,
    ".csv":     PARSE_TRACK_TABLE,
    ".log":     PARSE_TRACK_LOG,
    ".json":    PARSE_TRACK_LOG,
}


def get_parse_track(doc_file: Path) -> str:
    """根据文件后缀返回解析轨道名称"""
    return _TRACK_MAP.get(doc_file.suffix.lower(), PARSE_TRACK_RICH_TEXT)


# ── 通用工具 ───────────────────────────────────────────────────────────────────

def text_to_chunks(text: str) -> List[str]:
    """将纯文本按空行拆分为段落块"""
    return [chunk.strip() for chunk in text.split("\n\n") if chunk.strip()]


# ── 格式分派入口 ───────────────────────────────────────────────────────────────

def split_into_chunks(doc_file: Path) -> List[str]:
    """根据文件后缀分派到对应解析器，返回文本块列表"""
    suffix = doc_file.suffix.lower()
    # Track 1: 结构化文档
    if suffix in (".md", ".txt"):
        return _split_text_chunks(doc_file)
    if suffix in (".dita", ".ditamap", ".xml"):
        return _split_dita_chunks(doc_file)
    # Track 2: 富文本文档
    if suffix == ".docx":
        return _split_docx_chunks(doc_file)
    if suffix == ".pdf":
        return _split_pdf_chunks(doc_file)
    if suffix in (".ppt", ".pptx"):
        return _split_pptx_chunks(doc_file)
    # Track 3: 表格数据
    if suffix in (".xlsx", ".xls"):
        return _split_xlsx_chunks(doc_file)
    if suffix == ".csv":
        return _split_csv_chunks(doc_file)
    # Track 4: 日志/结构化记录
    if suffix == ".log":
        return _split_log_chunks(doc_file)
    if suffix == ".json":
        return _split_json_chunks(doc_file)
    # 未知格式：尝试当作纯文本处理
    try:
        return _split_text_chunks(doc_file)
    except UnicodeDecodeError:
        return []


# ── 纯文本 / Markdown ─────────────────────────────────────────────────────────

# chunk 粒度控制常量
MAX_CHUNK_CHARS = 800   # 超过此长度的 chunk 用滑窗二次切分
MIN_CHUNK_CHARS = 60    # 短于此长度的 chunk 与下一块合并


def _split_markdown_chunks(text: str, overlap_paragraphs: int = 2) -> List[str]:
    """Markdown 结构感知分块，规则优先级：

    1. 遇到标题行（# / ## / ### 等）开启新 chunk，标题 + 其下内容归为一块。
    2. 表格（连续含 | 的行）作为整体不拆散，追加到当前 chunk。
    3. 列表（连续以 - / * / 数字. 开头的行）整体不拆散，追加到当前 chunk。
    4. 普通段落（空行分隔）正常追加到当前 chunk。
    5. 相邻 chunk 之间保留 overlap_paragraphs 个段落的重叠，
       避免关键参数落在边界被切断。
    """
    import re

    _RE_HEADING  = re.compile(r"^#{1,6}\s+\S")
    _RE_TABLE_ROW = re.compile(r"^\s*\|")
    _RE_LIST_ITEM = re.compile(r"^\s*(?:[-*+]|\d+\.)\s+\S")

    lines = text.splitlines()

    # ── 第一步：将原始行归组为「语义段落」─────────────────────────────────────
    # 每个语义段落是一个 (kind, text) 元组，kind ∈ {heading, table, list, para}
    paragraphs: list[tuple[str, str]] = []
    i = 0
    while i < len(lines):
        line = lines[i]

        # 空行：跳过
        if not line.strip():
            i += 1
            continue

        # 标题行
        if _RE_HEADING.match(line):
            paragraphs.append(("heading", line.rstrip()))
            i += 1
            continue

        # 表格块：连续含 | 的行聚为一块
        if _RE_TABLE_ROW.match(line):
            block_lines = []
            while i < len(lines) and (_RE_TABLE_ROW.match(lines[i]) or not lines[i].strip()):
                if lines[i].strip():
                    block_lines.append(lines[i].rstrip())
                i += 1
            paragraphs.append(("table", "\n".join(block_lines)))
            continue

        # 列表块：连续列表项聚为一块（允许中间有缩进行）
        if _RE_LIST_ITEM.match(line):
            block_lines = []
            while i < len(lines):
                cur = lines[i]
                if not cur.strip():        # 空行结束列表
                    i += 1
                    break
                if _RE_LIST_ITEM.match(cur) or cur[:1] in (" ", "\t"):
                    block_lines.append(cur.rstrip())
                    i += 1
                else:
                    break
            paragraphs.append(("list", "\n".join(block_lines)))
            continue

        # 普通段落：遇到空行或不同类型行结束
        block_lines = []
        while i < len(lines):
            cur = lines[i]
            if not cur.strip():
                i += 1
                break
            if _RE_HEADING.match(cur) or _RE_TABLE_ROW.match(cur) or _RE_LIST_ITEM.match(cur):
                break
            block_lines.append(cur.rstrip())
            i += 1
        paragraphs.append(("para", "\n".join(block_lines)))

    # ── 第二步：将语义段落按标题分组为 chunk ──────────────────────────────────
    chunks_raw: list[list[tuple[str, str]]] = []   # 每个 chunk 是一组段落
    current: list[tuple[str, str]] = []

    for kind, body in paragraphs:
        if kind == "heading":
            # 遇到新标题：保存上一组（非空），开启新组
            if current:
                chunks_raw.append(current)
            current = [(kind, body)]
        else:
            current.append((kind, body))

    if current:
        chunks_raw.append(current)

    # ── 第三步：序列化为字符串，并加入段落级重叠 ──────────────────────────────
    def _serialize(paras: list[tuple[str, str]]) -> str:
        return "\n\n".join(body for _, body in paras if body.strip())

    result: list[str] = []
    for idx, group in enumerate(chunks_raw):
        # 前向重叠：把上一组的最后 overlap_paragraphs 个非标题段落附到本组开头
        prefix: list[tuple[str, str]] = []
        if overlap_paragraphs > 0 and idx > 0:
            prev_non_heading = [p for p in chunks_raw[idx - 1] if p[0] != "heading"]
            prefix = prev_non_heading[-overlap_paragraphs:]

        text_out = _serialize(prefix + group)
        if text_out.strip():
            result.append(text_out.strip())

    # ── 第四步：超长 chunk 按滑窗再切分，过短 chunk 与下一块合并 ─────────────
    split_result: list[str] = []
    for chunk in result:
        if len(chunk) <= MAX_CHUNK_CHARS:
            split_result.append(chunk)
            continue
        # 按段落分割后滑窗合并（window=4段落, step=2）
        paras = [p.strip() for p in chunk.split("\n\n") if p.strip()]
        window, step = 4, 2
        if len(paras) <= window:
            split_result.append(chunk)   # 拆不了则保留原块
            continue
        pos = 0
        while pos < len(paras):
            sub = "\n\n".join(paras[pos: pos + window])
            if sub.strip():
                split_result.append(sub.strip())
            pos += step

    # 合并过短 chunk
    merged: list[str] = []
    i = 0
    while i < len(split_result):
        if i < len(split_result) - 1 and len(split_result[i]) < MIN_CHUNK_CHARS:
            merged.append(split_result[i] + "\n\n" + split_result[i + 1])
            i += 2
        else:
            merged.append(split_result[i])
            i += 1

    return merged if merged else [text.strip()]


def _split_text_chunks(doc_file: Path) -> List[str]:
    content = doc_file.read_text(encoding="utf-8")
    if doc_file.suffix.lower() == ".md":
        return _split_markdown_chunks(content)
    return text_to_chunks(content)


# ── Word (.docx) ──────────────────────────────────────────────────────────────

def _split_docx_chunks(doc_file: Path) -> List[str]:
    """解析 Word (.docx) 文件，将段落合并为文本块"""
    try:
        from docx import Document  # type: ignore
    except ImportError as exc:
        raise ImportError("请先安装 python-docx：uv add python-docx") from exc

    doc = Document(str(doc_file))
    chunks: List[str] = []
    buffer: List[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            if buffer:
                chunks.append(" ".join(buffer))
                buffer = []
            continue
        if para.style and para.style.name.startswith("Heading"):
            if buffer:
                chunks.append(" ".join(buffer))
                buffer = []
            chunks.append(text)
        else:
            buffer.append(text)

    if buffer:
        chunks.append(" ".join(buffer))

    return [c for c in chunks if c]


# ── DITA / XML ────────────────────────────────────────────────────────────────

_DITA_TEXT_TAGS = {
    "title", "shortdesc", "p", "li", "note",
    "codeblock", "codeph", "term", "dd", "dt",
    "section", "result", "context", "steps",
}


def _iter_dita_text(element: ET.Element) -> str:
    """递归提取 DITA/XML 元素的可读文本，去除命名空间前缀和多余空白"""
    tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag
    parts: List[str] = []
    if tag in _DITA_TEXT_TAGS:
        raw = ET.tostring(element, encoding="unicode", method="text")
        text = " ".join(raw.split())
        if text:
            parts.append(text)
    else:
        for child in element:
            child_text = _iter_dita_text(child)
            if child_text:
                parts.append(child_text)
    return "\n".join(parts)


def extract_dita_metadata(doc_file: Path) -> dict[str, str]:
    """从 DITA/XML 文件的 prolog/metadata 节点自动抽取元数据。

    支持的 DITA 属性:
    - <prodinfo><prodname> → product_line
    - <prodinfo><vrmlist><vrm version="..."> → version
    - <othermeta name="department" content="..."> → department
    - <othermeta name="doc_id" content="..."> → doc_id
    - <othermeta name="owner" content="..."> → owner
    - <othermeta name="doc_type" content="..."> → doc_type
    - <othermeta name="status" content="..."> → status
    """
    from .config import META_FIELDS

    result: dict[str, str] = {}
    try:
        tree = ET.parse(str(doc_file))
        root = tree.getroot()
        for el in root.iter():
            tag = el.tag.split("}")[-1] if "}" in el.tag else el.tag
            if tag == "prodname" and el.text:
                result["product_line"] = el.text.strip()
            elif tag == "vrm":
                ver = el.get("version") or el.get("release", "")
                if ver:
                    result["version"] = ver.strip()
            elif tag == "othermeta":
                name = el.get("name", "")
                content = el.get("content", "")
                if name in META_FIELDS and content:
                    result[name] = content.strip()
    except Exception:
        pass
    return result


def _split_dita_chunks(doc_file: Path) -> List[str]:
    """解析 DITA / XML 文件，将各顶级元素的文本提取为块"""
    try:
        tree = ET.parse(str(doc_file))
    except ET.ParseError:
        try:
            return _split_text_chunks(doc_file)
        except UnicodeDecodeError:
            return []

    root = tree.getroot()
    chunks: List[str] = []

    title_el = root.find("title") or root.find("{*}title")
    if title_el is not None:
        title_text = " ".join((title_el.text or "").split())
        if title_text:
            chunks.append(title_text)

    for child in root:
        child_tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if child_tag == "title":
            continue
        text = _iter_dita_text(child).strip()
        if text:
            chunks.extend(text_to_chunks(text))

    return [c for c in chunks if c]


# ── PDF ───────────────────────────────────────────────────────────────────────

def _split_pdf_chunks(doc_file: Path) -> List[str]:
    """解析 PDF：文字版直接提取文字层；扫描件（文字层稀少）回退 Tesseract OCR"""
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise ImportError("请先安装 pymupdf：uv add pymupdf") from exc

    page_texts: List[str] = []
    total_chars = 0
    with fitz.open(str(doc_file)) as pdf:
        num_pages = max(len(pdf), 1)
        for page in pdf:
            text = page.get_text("text").strip()
            page_texts.append(text)
            total_chars += len(text)

    # 平均每页字符数 >= 100 视为文字版 PDF
    if total_chars / num_pages >= 100:
        chunks: List[str] = []
        for text in page_texts:
            if text:
                chunks.extend(text_to_chunks(text))
        return chunks

    # 扫描件：回退 Poppler + Tesseract OCR
    return _unstructured_pdf_chunks(doc_file)


def _unstructured_pdf_chunks(doc_file: Path) -> List[str]:
    """使用 Poppler + Tesseract OCR 解析扫描件 PDF（每页转图片后识别）"""
    import shutil
    import subprocess
    import tempfile

    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        raise RuntimeError("未找到 pdftoppm，请先安装 Poppler：brew install poppler")

    try:
        import unstructured_pytesseract as pytesseract  # type: ignore
        from PIL import Image  # type: ignore
    except ImportError as exc:
        raise ImportError("请先安装依赖：uv add 'unstructured[pdf]'") from exc

    chunks: List[str] = []
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            [pdftoppm, "-png", "-r", "200", str(doc_file), f"{tmp}/page"],
            check=True, capture_output=True,
        )
        for page_img in sorted(Path(tmp).glob("page-*.png")):
            img = Image.open(page_img)
            text = pytesseract.image_to_string(img, lang="chi_sim+eng").strip()
            if text:
                chunks.extend(text_to_chunks(text))
    return chunks


# ── PowerPoint (.pptx) ────────────────────────────────────────────────────────

def _split_pptx_chunks(doc_file: Path) -> List[str]:
    """解析 PPTX 文件，每张幻灯片提取标题和文本框内容合并为一个块"""
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Pt  # noqa: F401
    except ImportError as exc:
        raise ImportError("请先安装 python-pptx：uv add python-pptx") from exc

    prs = Presentation(str(doc_file))
    chunks: List[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    parts.append(line)
        if parts:
            chunks.append(f"[幻灯片 {slide_idx}] " + " ".join(parts))
    return chunks


# ── Track 3: 表格数据 (.xlsx / .xls / .csv) ──────────────────────────────────

def _split_xlsx_chunks(doc_file: Path) -> List[str]:
    """解析 Excel 文件：每个 Sheet 的每行转为"字段: 值 | 字段: 值"文本段"""
    try:
        import pandas as pd  # type: ignore
    except ImportError as exc:
        raise ImportError("请先安装依赖：uv add pandas openpyxl") from exc

    chunks: List[str] = []
    xf = pd.ExcelFile(str(doc_file))
    for sheet_name in xf.sheet_names:
        df = xf.parse(sheet_name).fillna("")
        headers = " | ".join(str(col) for col in df.columns)
        chunks.append(f"[Sheet: {sheet_name}] 字段: {headers}")
        for _, row in df.iterrows():
            row_text = " | ".join(
                f"{col}: {val}"
                for col, val in zip(df.columns, row)
                if str(val).strip()
            )
            if row_text.strip():
                chunks.append(row_text)
    return [c for c in chunks if c]


def _split_csv_chunks(doc_file: Path) -> List[str]:
    """解析 CSV 文件：每行转为"字段: 值 | 字段: 值"文本段"""
    try:
        import pandas as pd  # type: ignore
    except ImportError as exc:
        raise ImportError("请先安装依赖：uv add pandas") from exc

    df = pd.read_csv(str(doc_file), encoding="utf-8-sig").fillna("")
    headers = " | ".join(str(col) for col in df.columns)
    chunks: List[str] = [f"字段: {headers}"]
    for _, row in df.iterrows():
        row_text = " | ".join(
            f"{col}: {val}"
            for col, val in zip(df.columns, row)
            if str(val).strip()
        )
        if row_text.strip():
            chunks.append(row_text)
    return [c for c in chunks if c]


# ── Track 4: 日志/结构化记录 (.log / .json) ──────────────────────────────────

_LOG_CHUNK_LINES = 30  # 每块聚合行数


def _split_log_chunks(doc_file: Path) -> List[str]:
    """解析日志文件：将连续日志行按固定窗口聚合为块，保留 ERROR/WARN 上下文"""
    content = doc_file.read_text(encoding="utf-8", errors="replace")
    lines = [l for l in content.splitlines() if l.strip()]
    chunks: List[str] = []
    for i in range(0, len(lines), _LOG_CHUNK_LINES):
        block = "\n".join(lines[i : i + _LOG_CHUNK_LINES])
        chunks.append(block)
    return chunks


def _split_json_chunks(doc_file: Path) -> List[str]:
    """解析 JSON 文件：顶层数组→每个对象为一段；顶层对象→每个 key 为一段"""
    content = doc_file.read_text(encoding="utf-8")
    try:
        data = json.loads(content)
    except json.JSONDecodeError:
        return text_to_chunks(content)

    chunks: List[str] = []
    if isinstance(data, list):
        for item in data:
            if isinstance(item, dict):
                text = " | ".join(
                    f"{k}: {v}" for k, v in item.items() if v is not None and str(v).strip()
                )
                if text:
                    chunks.append(text)
            else:
                chunks.append(str(item))
    elif isinstance(data, dict):
        for k, v in data.items():
            chunks.append(f"{k}: {json.dumps(v, ensure_ascii=False)}")
    else:
        chunks.append(str(data))
    return [c for c in chunks if c]
